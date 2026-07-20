from __future__ import annotations

from collections.abc import Iterable
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ingestion.models import SourceUnit
from ingestion.utils import (
    clean_text,
    relative_path,
    safe_filename,
)


def _walk_shapes(
    shapes: Iterable[Any],
) -> Iterable[Any]:
    """
    슬라이드의 모든 shape를 순회한다.

    Group shape 내부에 들어 있는 텍스트와 이미지도
    놓치지 않도록 재귀적으로 처리한다.
    """

    ordered_shapes = sorted(
        shapes,
        key=lambda shape: (
            int(getattr(shape, "top", 0)),
            int(getattr(shape, "left", 0)),
        ),
    )

    for shape in ordered_shapes:
        if (
            shape.shape_type
            == MSO_SHAPE_TYPE.GROUP
        ):
            yield from _walk_shapes(
                shape.shapes
            )
            continue

        yield shape


def _extract_table_text(
    shape: Any,
) -> str:
    rows: list[str] = []

    for row in shape.table.rows:
        cell_values = [
            clean_text(cell.text)
            for cell in row.cells
        ]

        if any(cell_values):
            rows.append(
                " | ".join(cell_values)
            )

    return "\n".join(rows)


def _extract_text_frame(
    shape: Any,
) -> str:
    lines: list[str] = []

    for paragraph in (
        shape.text_frame.paragraphs
    ):
        paragraph_text = clean_text(
            paragraph.text
        )

        if not paragraph_text:
            continue

        indentation = (
            "  " * paragraph.level
        )

        lines.append(
            f"{indentation}{paragraph_text}"
        )

    return "\n".join(lines)


def _extract_shape_text(
    shape: Any,
) -> str:
    if getattr(
        shape,
        "has_table",
        False,
    ):
        return _extract_table_text(shape)

    if getattr(
        shape,
        "has_text_frame",
        False,
    ):
        return _extract_text_frame(shape)

    return ""


def _save_picture(
    shape: Any,
    output_directory: Path,
    project_root: Path,
    slide_number: int,
    image_number: int,
) -> str:
    extension = (
        shape.image.ext or "bin"
    ).lower()

    output_path = (
        output_directory
        / (
            f"slide-{slide_number:03d}-"
            f"image-{image_number:02d}."
            f"{extension}"
        )
    )

    output_path.write_bytes(
        shape.image.blob
    )

    return relative_path(
        output_path,
        project_root,
    )


def _get_slide_title(
    slide: Any,
    content: str,
    fallback: str,
) -> str:
    title_shape = slide.shapes.title

    if title_shape is not None:
        title_text = clean_text(
            title_shape.text
        )

        if title_text:
            return title_text[:140]

    content_lines = [
        line.strip()
        for line in content.splitlines()
        if line.strip()
    ]

    if content_lines:
        return content_lines[0][:140]

    return fallback


def parse_pptx(
    file_path: Path,
    image_directory: Path,
    project_root: Path,
) -> list[SourceUnit]:
    presentation = Presentation(
        file_path
    )

    source_file = relative_path(
        file_path,
        project_root,
    )

    file_prefix = safe_filename(
        file_path.stem
    )

    presentation_image_directory = (
        image_directory / file_prefix
    )

    presentation_image_directory.mkdir(
        parents=True,
        exist_ok=True,
    )

    units: list[SourceUnit] = []

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1,
    ):
        text_blocks: list[str] = []
        image_paths: list[str] = []

        image_number = 0

        for shape in _walk_shapes(
            slide.shapes
        ):
            shape_text = (
                _extract_shape_text(shape)
            )

            if (
                shape_text
                and shape_text
                not in text_blocks
            ):
                text_blocks.append(
                    shape_text
                )

            if (
                shape.shape_type
                == MSO_SHAPE_TYPE.PICTURE
            ):
                image_number += 1

                saved_image_path = (
                    _save_picture(
                        shape=shape,
                        output_directory=(
                            presentation_image_directory
                        ),
                        project_root=(
                            project_root
                        ),
                        slide_number=(
                            slide_number
                        ),
                        image_number=(
                            image_number
                        ),
                    )
                )

                image_paths.append(
                    saved_image_path
                )

        content = clean_text(
            "\n".join(text_blocks)
        )

        fallback_title = (
            f"{file_path.stem} "
            f"— Slide {slide_number}"
        )

        title = _get_slide_title(
            slide=slide,
            content=content,
            fallback=fallback_title,
        )

        unit = SourceUnit(
            title=title,
            content=content,
            source_file=source_file,
            file_type="pptx",
            section_path=[title],
            slide_number=slide_number,
            image_paths=image_paths,
            kind="slide",
            ocr_applied=False,
            requires_visual_review=(
                bool(image_paths)
                or not content
            ),
            should_display_image=(
                bool(image_paths)
            ),
            metadata={
                "slide_number": (
                    slide_number
                ),
                "image_count": len(
                    image_paths
                ),
                "text_character_count": (
                    len(content)
                ),
            },
        )

        # 슬라이드마다 정확히 한 번만 추가한다.
        units.append(unit)

    return units